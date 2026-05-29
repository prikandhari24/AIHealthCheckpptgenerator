from __future__ import annotations

from functools import lru_cache
import re

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

from src.config import (
    CALLOUT_BORDER_COLOR,
    CALLOUT_FILL_COLOR,
    CALLOUT_FONT_SIZE,
    CALLOUT_HEIGHT_FACTOR,
    CALLOUT_RIGHT_PADDING,
    CALLOUT_TEXT_COLOR,
    CALLOUT_WIDTH,
    HEADER_CLIENT_OFFSET,
    HEADER_PEER_OFFSET,
    LEGEND_FONT_SIZE,
    LEGEND_ITEM_GAP,
    LEGEND_LINE_GAP,
    LEGEND_LOGO_SCALE,
    LEGEND_MAX_LOGO_HEIGHT,
    LEGEND_RIGHT_MARGIN,
    LEGEND_TEXT_GAP,
    OUTPUT_DIR,
    OUTPUT_PATH,
    ROW_LOGO_GAP,
    ROW_LOGO_HORIZONTAL_PADDING,
    ROW_LOGO_MAX_HEIGHT_FACTOR,
    ROW_LOGO_VERTICAL_GAP,
    SCORE_SLIDE_INDEXES,
    TEMPLATE_PATH,
)
from src.excel_parser import normalize_label

LEVEL_RANK = {
    "behind": 0,
    "at_par": 1,
    "ahead": 2,
}
MAX_CALLOUTS_PER_SLIDE = 3
MAX_CALLOUT_TEXT_LENGTH = 70


STATUS_TO_COLUMN = {
    "behind": 3,
    "at_par": 4,
    "ahead": 5,
}

CLIENT_NAME_PLACEHOLDERS = ("Legalitas", "Legálitas")


def _find_text_shape(slide, target_text: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            if shape.text.strip() == target_text:
                return shape
    return None


def _find_table_shape(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape
    return None


def _replace_text_in_slide(slide, replacements: dict[str, str]):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                updated_text = run.text
                for source, target in replacements.items():
                    updated_text = updated_text.replace(source, target)
                run.text = updated_text


def _replace_client_name_across_deck(presentation, client_company: str):
    replacements = {
        placeholder: client_company
        for placeholder in CLIENT_NAME_PLACEHOLDERS
    }
    for slide in presentation.slides:
        _replace_text_in_slide(slide, replacements)


def _replace_title_slide_logo(slide, client_logo_path: str):
    picture_shapes = [
        shape for shape in slide.shapes
        if shape.shape_type == 13
    ]
    if not picture_shapes:
        return

    target_shape = picture_shapes[0]
    left = target_shape.left
    top = target_shape.top
    width = target_shape.width
    height = target_shape.height

    target_shape._element.getparent().remove(target_shape._element)
    slide.shapes.add_picture(
        client_logo_path,
        left,
        top,
        width=width,
        height=height,
    )


@lru_cache(maxsize=None)
def _logo_aspect_ratio(logo_path: str) -> float:
    with Image.open(logo_path) as image:
        return image.width / max(image.height, 1)


def _cell_bounds(table_shape, row_idx: int, col_idx: int) -> tuple[int, int, int, int]:
    table = table_shape.table
    left = table_shape.left + sum(table.columns[index].width for index in range(col_idx))
    top = table_shape.top + sum(table.rows[index].height for index in range(row_idx))
    width = table.columns[col_idx].width
    height = table.rows[row_idx].height
    return left, top, width, height


def _row_logo_layout_values(width: int, height: int) -> tuple[int, int, int, int]:
    padding_x = min(int(ROW_LOGO_HORIZONTAL_PADDING), max(width // 10, 1))
    gap_x = min(int(ROW_LOGO_GAP), max(width // 40, 1))
    gap_y = min(int(ROW_LOGO_VERTICAL_GAP), max(height // 20, 1))
    max_height = max(int(height * ROW_LOGO_MAX_HEIGHT_FACTOR), 1)
    return padding_x, gap_x, gap_y, max_height


def _single_row_layout_height(width: int, height: int, logo_paths: list[str]) -> int:
    if not logo_paths:
        return 0

    padding_x, gap_x, gap_y, max_height = _row_logo_layout_values(width, height)
    available_width = max(width - (2 * padding_x), 1)
    aspect_sum = sum(_logo_aspect_ratio(path) for path in logo_paths)
    width_budget = max(available_width - (gap_x * (len(logo_paths) - 1)), 1)
    return max(1, min(max_height, int(width_budget / max(aspect_sum, 0.01))))


def _compute_chart_logo_height(presentation, parsed_data, company_order: list[str], logo_paths_by_company: dict[str, str]) -> int:
    feasible_heights = []

    for slide_index in SCORE_SLIDE_INDEXES:
        if slide_index >= len(presentation.slides):
            continue

        table_shape = _find_table_shape(presentation.slides[slide_index])
        if table_shape is None:
            continue

        table = table_shape.table
        for row_idx in range(1, len(table.rows)):
            indicator = table.cell(row_idx, 1).text.strip()
            if not indicator:
                continue

            assessment = parsed_data.get_indicator_assessment(indicator)
            if assessment is None:
                continue

            for status, col_idx in STATUS_TO_COLUMN.items():
                matching_logos = [
                    logo_paths_by_company[company]
                    for company in company_order
                    if assessment.company_levels.get(company) == status and company in logo_paths_by_company
                ]
                if not matching_logos:
                    continue

                _, _, cell_width, cell_height = _cell_bounds(table_shape, row_idx, col_idx)
                candidate_height = _single_row_layout_height(
                    cell_width,
                    cell_height,
                    matching_logos,
                )
                feasible_heights.append(candidate_height)

    if not feasible_heights:
        return 1

    return max(1, min(feasible_heights))


def _add_logo_strip(
    slide,
    logo_paths: list[str],
    left: int,
    top: int,
    width: int,
    height: int,
    render_height: int,
):
    if not logo_paths:
        return

    padding_x, gap_x, gap_y, max_height = _row_logo_layout_values(width, height)
    bounded_height = max(1, min(render_height, max_height))
    render_widths = [
        max(int(bounded_height * _logo_aspect_ratio(path)), 1)
        for path in logo_paths
    ]
    total_width = sum(render_widths) + gap_x * (len(render_widths) - 1)
    start_left = left + max((width - total_width) // 2, 0)
    start_top = top + max((height - bounded_height) // 2, 0)

    cursor_left = start_left
    for logo_path, render_width in zip(logo_paths, render_widths):
        slide.shapes.add_picture(
            logo_path,
            cursor_left,
            start_top,
            height=bounded_height,
        )
        cursor_left += render_width + gap_x


def _estimate_text_width(text: str) -> int:
    font_points = LEGEND_FONT_SIZE.pt
    # Approximate average character width for compact sans-serif legend text.
    return int(len(text) * font_points * 7000) + int(LEGEND_TEXT_GAP)


def _add_callout_box(slide, text: str, left: int, top: int, width: int, height: int):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CALLOUT_FILL_COLOR
    shape.line.color.rgb = CALLOUT_BORDER_COLOR

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = CALLOUT_FONT_SIZE
    run.font.color.rgb = CALLOUT_TEXT_COLOR


def _add_legend_text(slide, text: str, left: int, top: int, width: int, height: int):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = False
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = LEGEND_FONT_SIZE


def _draw_legend_item(slide, company: str, logo_path: str, left: int, top: int, logo_height: int) -> int:
    logo_width = max(int(logo_height * _logo_aspect_ratio(logo_path)), 1)
    text_width = max(_estimate_text_width(company), 1)

    slide.shapes.add_picture(
        logo_path,
        left,
        top,
        height=logo_height,
    )
    text_left = left + logo_width + int(LEGEND_TEXT_GAP)
    _add_legend_text(
        slide,
        company,
        text_left,
        top,
        text_width,
        logo_height,
    )
    return logo_width + int(LEGEND_TEXT_GAP) + text_width


def _summary_text_for_indicator(parsed_data, client_company: str, indicator: str) -> str:
    summary_column = parsed_data.summary_columns.get(client_company)
    if not summary_column:
        return ""

    indicator_key = normalize_label(indicator)
    matching_rows = parsed_data.raw_df[
        parsed_data.raw_df["Indicator"].astype(str).map(normalize_label) == indicator_key
    ]
    if matching_rows.empty:
        return ""

    value = " ".join(str(matching_rows.iloc[0][summary_column]).split()).strip()
    if not value or value.lower() == "nan":
        return ""
    return _callout_reason_text(value)


def _callout_reason_text(summary_text: str) -> str:
    text = " ".join(summary_text.split()).strip()
    if not text:
        return ""

    but_match = re.search(r"\bbut\b", text, flags=re.IGNORECASE)
    if but_match:
        text = text[but_match.end():].strip()

    text = re.sub(r"^assessed\s+level\s*:\s*", "", text, flags=re.IGNORECASE).strip()
    text = re.split(r"[.;]", text, maxsplit=1)[0].strip()
    text = re.sub(r"^[,\-: ]+", "", text).strip()
    text = text.rstrip(". ")

    if len(text) > MAX_CALLOUT_TEXT_LENGTH:
        truncated = text[:MAX_CALLOUT_TEXT_LENGTH].rsplit(" ", 1)[0].strip()
        text = (truncated or text[:MAX_CALLOUT_TEXT_LENGTH]).rstrip(",;:- ") + "..."

    return text


def _peer_advantage_count(assessment, client_company: str, peer_companies: list[str]) -> int:
    client_level = assessment.company_levels.get(client_company)
    if client_level not in LEVEL_RANK or not peer_companies:
        return 0

    client_rank = LEVEL_RANK[client_level]
    return sum(
        1
        for company in peer_companies
        if LEVEL_RANK.get(assessment.company_levels.get(company), client_rank) > client_rank
    )


def _select_slide_callouts(table, parsed_data, client_company: str, peer_companies: list[str]) -> dict[int, str]:
    candidates: list[tuple[int, int, str]] = []

    for row_idx in range(1, len(table.rows)):
        indicator = table.cell(row_idx, 1).text.strip()
        if not indicator:
            continue

        assessment = parsed_data.get_indicator_assessment(indicator)
        if assessment is None:
            continue

        better_peer_count = _peer_advantage_count(
            assessment,
            client_company,
            peer_companies,
        )
        if better_peer_count <= 0:
            continue

        summary_text = _summary_text_for_indicator(
            parsed_data,
            client_company,
            assessment.indicator,
        )
        if not summary_text:
            continue

        candidates.append((row_idx, better_peer_count, summary_text))

    candidates.sort(key=lambda candidate: (-candidate[1], candidate[0]))
    return {
        row_idx: summary_text
        for row_idx, _, summary_text in candidates[:MAX_CALLOUTS_PER_SLIDE]
    }


def _render_callout_for_row(slide, table_shape, row_idx: int, callout_text: str):
    desc_left, desc_top, desc_width, desc_height = _cell_bounds(table_shape, row_idx, 2)
    callout_width = min(int(CALLOUT_WIDTH), max(desc_width // 3, 1))
    callout_height = max(1, int(desc_height * CALLOUT_HEIGHT_FACTOR))
    callout_left = desc_left + desc_width - callout_width - int(CALLOUT_RIGHT_PADDING)
    callout_top = desc_top + max((desc_height - callout_height) // 2, 0)
    _add_callout_box(
        slide,
        callout_text,
        callout_left,
        callout_top,
        callout_width,
        callout_height,
    )


def _add_header_legends(
    slide,
    slide_width: int,
    client_company: str,
    client_logo_path: str,
    peer_companies: list[str],
    logo_paths_by_company: dict[str, str],
    chart_logo_height: int,
):
    client_label = _find_text_shape(slide, "Client:")
    peer_label = _find_text_shape(slide, "Peers:")
    legend_logo_height = max(
        1,
        min(
            int(chart_logo_height * LEGEND_LOGO_SCALE),
            int(LEGEND_MAX_LOGO_HEIGHT),
        ),
    )

    if client_label is not None:
        client_top = client_label.top + max((client_label.height - legend_logo_height) // 2, 0)
        _draw_legend_item(
            slide,
            client_company,
            client_logo_path,
            client_label.left + int(HEADER_CLIENT_OFFSET),
            client_top,
            legend_logo_height,
        )

    if peer_label is not None:
        peer_top = peer_label.top + max((peer_label.height - legend_logo_height) // 2, 0)
        cursor_left = peer_label.left + int(HEADER_PEER_OFFSET)
        cursor_top = peer_top
        max_right = slide_width - int(LEGEND_RIGHT_MARGIN)

        for company in peer_companies:
            if company not in logo_paths_by_company:
                continue

            logo_path = logo_paths_by_company[company]
            item_width = (
                max(int(legend_logo_height * _logo_aspect_ratio(logo_path)), 1)
                + int(LEGEND_TEXT_GAP)
                + _estimate_text_width(company)
            )
            if cursor_left + item_width > max_right and cursor_left > peer_label.left + int(HEADER_PEER_OFFSET):
                cursor_left = peer_label.left + int(HEADER_PEER_OFFSET)
                cursor_top += legend_logo_height + int(LEGEND_LINE_GAP)

            used_width = _draw_legend_item(
                slide,
                company,
                logo_path,
                cursor_left,
                cursor_top,
                legend_logo_height,
            )
            cursor_left += used_width + int(LEGEND_ITEM_GAP)


def _render_slide_scores(
    slide,
    parsed_data,
    client_company: str,
    peer_companies: list[str],
    company_order: list[str],
    logo_paths_by_company: dict[str, str],
    chart_logo_height: int,
    include_callouts: bool,
):
    table_shape = _find_table_shape(slide)
    if table_shape is None:
        raise ValueError("No table found on score slide.")

    table = table_shape.table
    missing_indicators = []
    slide_callouts = (
        _select_slide_callouts(
            table,
            parsed_data,
            client_company,
            peer_companies,
        )
        if include_callouts
        else {}
    )

    for row_idx in range(1, len(table.rows)):
        indicator = table.cell(row_idx, 1).text.strip()
        if not indicator:
            continue

        assessment = parsed_data.get_indicator_assessment(indicator)
        if assessment is None:
            missing_indicators.append(indicator)
            continue

        for status, col_idx in STATUS_TO_COLUMN.items():
            matching_logos = [
                logo_paths_by_company[company]
                for company in company_order
                if assessment.company_levels.get(company) == status and company in logo_paths_by_company
            ]
            cell_left, cell_top, cell_width, cell_height = _cell_bounds(table_shape, row_idx, col_idx)
            _add_logo_strip(
                slide,
                matching_logos,
                cell_left,
                cell_top,
                cell_width,
                cell_height,
                chart_logo_height,
            )

        if row_idx in slide_callouts:
            _render_callout_for_row(
                slide,
                table_shape,
                row_idx,
                slide_callouts[row_idx],
            )

    if missing_indicators:
        raise ValueError(
            "These template indicators were not found in the Excel file: "
            + ", ".join(missing_indicators)
        )


def build_ppt(
    parsed_data,
    client_company: str,
    title_logo_path: str,
    client_logo_path: str,
    peer_logo_paths: dict[str, str],
    include_callouts: bool = True,
):
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    presentation = Presentation(str(TEMPLATE_PATH))
    _replace_client_name_across_deck(presentation, client_company)
    if presentation.slides:
        _replace_title_slide_logo(presentation.slides[0], title_logo_path)

    ordered_companies = [client_company] + [
        company for company in parsed_data.companies if company != client_company
    ]
    peer_companies = [
        company for company in parsed_data.companies if company != client_company
    ]
    logo_paths_by_company = {client_company: client_logo_path, **peer_logo_paths}
    chart_logo_height = _compute_chart_logo_height(
        presentation,
        parsed_data,
        ordered_companies,
        logo_paths_by_company,
    )

    for slide_index in SCORE_SLIDE_INDEXES:
        if slide_index >= len(presentation.slides):
            continue

        slide = presentation.slides[slide_index]
        _add_header_legends(
            slide,
            presentation.slide_width,
            client_company,
            client_logo_path,
            peer_companies,
            logo_paths_by_company,
            chart_logo_height,
        )
        _render_slide_scores(
            slide,
            parsed_data,
            client_company,
            peer_companies,
            ordered_companies,
            logo_paths_by_company,
            chart_logo_height,
            include_callouts,
        )

    OUTPUT_DIR.mkdir(exist_ok=True)
    presentation.save(str(OUTPUT_PATH))
    return str(OUTPUT_PATH)
